from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.contrib.auth.decorators import login_required

from certificate.reupload import reupload

from .models import Event, Participant
import pandas as pd
from .upload import  upload
from pptx import Presentation
from django.core.mail import send_mail, EmailMessage
import requests
import os
import webbrowser


def verify(request):
	if request.method=='POST':
		certId=request.POST.get('cert_id')
		url ="https://docs.google.com/presentation/d/"+certId+"/export/pdf"
		if os.path.exists("mycertificate.pdf"):
			os.remove("mycertificate.pdf")
		x = requests.head(url,allow_redirects=True)
		if x.headers.get('Content-Type')=="application/pdf":
			f=requests.get(url,allow_redirects=True)
			open("mycertificate.pdf",'wb').write(f.content)
			with open("mycertificate.pdf", 'rb') as fh:
				response = HttpResponse(fh.read(), content_type="application/pdf")
				response['Content-Disposition'] = 'inline; filename=' + os.path.basename("mycertificate.pdf")
				return response
		else:
			return render(request,"invalid.html")
	else:
		return redirect("/")

def index(request):
	return render(request, 'index.html')

@login_required
def create(request):
	if request.method == "POST":
		csv = request.FILES.get('csv')
		temp = request.FILES.get('template')
		event_name = request.POST.get('event_name')
		
		event = Event(user = request.user,
			event_name = event_name,
			csv_file = csv,
			template = temp)
		event.save()

		return redirect(f"/certificate/{event.id}/{event.slug}")

	return render(request, 'certificate/create_event.html')

@login_required
def delete_event(request, id, slug):
	event = Event.objects.filter(slug=slug, id=id).first()
	if event.user == request.user:
		event.delete()
	return redirect('view_certificate_status')

@login_required
def track(request, id, slug):
	event = Event.objects.filter(slug=slug, id=id).first()
	if event.message:

		return render(request, 'certificate/track.html', {
			'event_name': event.event_name,
			'event_date': event.date,
			'participat_details': Participant.objects.filter(event=event)
			})

	prs = Presentation(event.template)
	st=""
	for slide in prs.slides:
		for shape in slide.shapes:
			if shape.has_text_frame:
				st = st + shape.text
				st = st + " "
	li = st.split()
	tags = []
	for i in li:
		if i[0] == "<" and i[-1] == ">":
			tags.append(i)

	if request.method == "POST":
		email_col = request.POST.get('emails')
		subject = request.POST.get('subject')
		mess = request.POST.get('mess')
		values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]	
		
		event.email_column = email_col
		event.message = mess
		event.subject = subject
		event.save()

		df=pd.read_csv(event.csv_file)
		df_len=df.shape
		i=0
	
		



		li=["First","Second","Third"]
		while i < df_len[0]:
			prs = Presentation(event.template)
			s_name = df.loc[i,event.email_column].split('@')[0]
			c_id=upload(s_name,"sample.pptx")
			
			for tag, v_type, value in values:
				for slide in prs.slides:
					for shape in slide.shapes:
						if shape.has_text_frame:
							if(shape.text.find(tag))!=-1:
									text_frame = shape.text_frame
									for paragraph in text_frame.paragraphs:
										for run in paragraph.runs:
											cur_text = run.text
											if v_type == 'text':
												new_text = cur_text.replace(tag, value)
											elif v_type == 'date':
												new_text = cur_text.replace(tag, '/'.join(value.split('-')[::-1]))
											elif v_type == 'csv':
												new_text = cur_text.replace(tag, df.loc[i,value])
											elif v_type == "auto":
												new_text = cur_text.replace(tag,c_id)
											else:
												pass
											run.text = new_text
											
			prs.save(s_name+".pptx")
			reupload(c_id,s_name+".pptx")
			url ="https://docs.google.com/presentation/d/"+c_id+"/export/pdf"
			f=requests.get(url,allow_redirects=True)
			open("certificate.pdf", 'wb').write(f.content)


			try:
				mail = EmailMessage(subject,
					f"Hello, {s_name} \n{mess}",
					settings.EMAIL_HOST_USER,
					[df.loc[i,event.email_column]])
				mail.attach_file("certificate.pdf")
				mail.send()
				Participant(event=event, email=df.loc[i,event.email_column], status=True).save()
				os.remove("certificate.pdf")
				os.remove(s_name+".pptx")
			except:
				Participant(event=event, email=df.loc[i,event.email_column], status=False).save()
				os.remove("certificate.pdf")
				os.remove(s_name+".pptx")
			i=i+1

		messages.success(request, "Certificates Sent Successfuly !!")
		return redirect(f"/certificate/{event.id}/{event.slug}")


	return render(request, 'certificate/map_tags_of_template.html',{
		'columns': list(pd.read_csv(event.csv_file).columns),
		'tags': tags,
		})


@login_required
def view_certificate_status(request):
	return render(request, 'certificate/view_certificate_status.html',{
		'events': Event.objects.filter(user=request.user)
		})



